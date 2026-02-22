from __future__ import annotations

import json
import shutil
import subprocess
import tempfile
import uuid
from pathlib import Path
from typing import Any, Dict


GENERATOR_PATH = Path(__file__).resolve().parent / "generate_report.js"
REPORTS_DIR = Path("generated")


def _has_node() -> bool:
    return shutil.which("node") is not None


def generate_clinical_report_pptx(
    structured_data: Dict[str, Any],
    insights: Dict[str, Any] | None = None,
    session_id: str | None = None,
) -> Dict[str, Any] | None:
    """
    Generate a clinical PPTX using Node + PptxGenJS.
    Returns file metadata consumable by presenter/server download flow, or None on any failure.
    """
    if not structured_data or not isinstance(structured_data, dict):
        return None

    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    rid = str(uuid.uuid4())[:8]
    out_name = f"Clinical_Report_{session_id or rid}.pptx"
    out_path = REPORTS_DIR / out_name
    payload = {
        "structured_data": structured_data,
        "insights": insights or {},
    }

    temp_json = Path(tempfile.gettempdir()) / f"clinical_report_{rid}.json"
    try:
        temp_json.write_text(json.dumps(payload), encoding="utf-8")
        if GENERATOR_PATH.exists() and _has_node():
            run = subprocess.run(
                ["node", str(GENERATOR_PATH), str(temp_json), str(out_path)],
                capture_output=True,
                text=True,
                timeout=90,
            )
            if run.returncode == 0 and out_path.exists():
                return {
                    "type": "pptx",
                    "name": out_path.name,
                    "path": str(out_path),
                    "size": int(out_path.stat().st_size),
                }
        fallback = _generate_with_python_pptx(out_path, structured_data, insights or {})
        if fallback:
            return {
                "type": "pptx",
                "name": out_path.name,
                "path": str(out_path),
                "size": int(out_path.stat().st_size),
            }
        return None
    except Exception:
        return None
    finally:
        try:
            temp_json.unlink(missing_ok=True)
        except Exception:
            pass


def generate_clinical_report_docx(
    structured_data: Dict[str, Any],
    insights: Dict[str, Any] | None = None,
    session_id: str | None = None,
) -> Dict[str, Any] | None:
    """
    Generate a structured clinical DOCX report.
    Returns file metadata consumable by presenter/server download flow, or None on failure.
    """
    if not structured_data or not isinstance(structured_data, dict):
        return None

    try:
        from docx import Document  # type: ignore
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
        from docx.shared import Pt, RGBColor, Inches  # type: ignore
        from docx.oxml import OxmlElement  # type: ignore
        from docx.oxml.ns import qn  # type: ignore
    except Exception:
        return None

    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    rid = str(uuid.uuid4())[:8]
    out_name = f"Clinical_Report_{session_id or rid}.docx"
    out_path = REPORTS_DIR / out_name
    ins = insights or {}

    try:
        doc = Document()
        theme = {
            "title": "0B1020",
            "subtitle": "475569",
            "accent": "4CC9F0",
            "section": "1D4ED8",
            "body": "E2E8F0",
            "muted": "94A3B8",
            "danger": "FF4D6D",
            "warning": "FFB703",
            "success": "2AEF9B",
            "panel_dark": "0F172A",
            "panel_mid": "1E293B",
            "panel_soft": "334155",
        }
        section = doc.sections[0]
        section.top_margin = Inches(0.6)
        section.bottom_margin = Inches(0.6)
        section.left_margin = Inches(0.7)
        section.right_margin = Inches(0.7)

        def _shade_cell(cell: Any, color: str) -> None:
            tc_pr = cell._tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:fill"), color)
            tc_pr.append(shd)

        def _set_run_style(run: Any, *, size: int = 11, bold: bool = False, color: str = "0F172A") -> None:
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)

        def _add_h1(text: str) -> None:
            p = doc.add_paragraph()
            run = p.add_run(text)
            _set_run_style(run, size=15, bold=True, color=theme["section"])
            p.paragraph_format.space_before = Pt(10)
            p.paragraph_format.space_after = Pt(6)

        def _add_small(text: str, color: str = "64748B") -> None:
            p = doc.add_paragraph()
            run = p.add_run(text)
            _set_run_style(run, size=9, bold=False, color=color)
            p.paragraph_format.space_after = Pt(3)

        chief = ((structured_data.get("chief_complaint") or {}).get("complaint")) or "not captured"
        patient = structured_data.get("patient") or {}
        diagnosis = ins.get("diagnosis_support") or {}
        actionable = ins.get("actionable_insights") or {}
        analytics = ins.get("analytics") or {}

        conf = structured_data.get("confidence_score", 0.5)
        try:
            conf_pct = int(max(0.0, min(1.0, float(conf or 0.0))) * 100)
        except Exception:
            conf_pct = 50
        risk_tier = str(diagnosis.get("risk_tier") or "low")
        red_count = int((analytics.get("red_flag_summary") or {}).get("count") or 0)
        gap_count = len(actionable.get("documentation_gaps") or [])

        title = doc.add_paragraph()
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        t_run = title.add_run("Clinical Intelligence Report")
        _set_run_style(t_run, size=24, bold=True, color=theme["title"])
        subtitle = doc.add_paragraph()
        subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
        s_run = subtitle.add_run("Structured Decision Support Summary")
        _set_run_style(s_run, size=11, color=theme["subtitle"])
        _add_small(
            f"Session ID: {session_id or 'N/A'} | Generated by Cavista Clinical Assistant",
            color="64748B",
        )
        _add_small("Decision support only; clinician verification required.", color="B91C1C")

        kpi = doc.add_table(rows=2, cols=4)
        kpi.style = "Table Grid"
        headers = ["Risk Tier", "Red Flags", "Documentation Gaps", "Confidence"]
        values = [risk_tier.title(), str(red_count), str(gap_count), f"{conf_pct}%"]
        for i, h in enumerate(headers):
            h_cell = kpi.rows[0].cells[i]
            h_cell.text = h
            _shade_cell(h_cell, theme["panel_dark"])
            run = h_cell.paragraphs[0].runs[0]
            _set_run_style(run, size=10, bold=True, color=theme["accent"])
            v_cell = kpi.rows[1].cells[i]
            v_cell.text = values[i]
            _shade_cell(v_cell, theme["panel_mid"])
            v_run = v_cell.paragraphs[0].runs[0]
            if i == 0 and values[i].lower() == "high":
                v_color = theme["danger"]
            elif i == 0 and values[i].lower() == "medium":
                v_color = theme["warning"]
            elif i == 0 and values[i].lower() == "low":
                v_color = theme["success"]
            else:
                v_color = theme["body"]
            _set_run_style(v_run, size=12, bold=True, color=v_color)

        _add_h1("1. Executive Summary")
        p = doc.add_paragraph()
        run = p.add_run(f"Chief complaint: {chief}")
        _set_run_style(run, size=11, bold=True)
        p2 = doc.add_paragraph()
        p2.add_run(
            f"Patient profile: gender={patient.get('gender', 'unknown')}, age={patient.get('age', 'N/A')}"
        )
        p3 = doc.add_paragraph()
        p3.add_run(
            "This report is intended for clinical documentation support and prioritization, not final diagnosis."
        )

        _add_h1("2. Vitals Snapshot")
        vitals = analytics.get("vitals") or []
        if vitals:
            table = doc.add_table(rows=1, cols=5)
            table.style = "Table Grid"
            hdr = table.rows[0].cells
            for idx, txt in enumerate(["Vital", "Value", "Unit", "Normal Range", "Status"]):
                hdr[idx].text = txt
                _shade_cell(hdr[idx], theme["panel_dark"])
                _set_run_style(hdr[idx].paragraphs[0].runs[0], size=10, bold=True, color=theme["accent"])
            for v in vitals:
                row = table.add_row().cells
                row[0].text = str(v.get("name") or "")
                row[1].text = "not captured" if v.get("value") is None else str(v.get("value"))
                row[2].text = str(v.get("unit") or "")
                row[3].text = f"{v.get('min', '-')}" + " - " + f"{v.get('max', '-')}"
                row[4].text = str(v.get("status") or "not_captured")
                status = str(v.get("status") or "not_captured").lower()
                if status == "high":
                    _shade_cell(row[4], "3F1C2B")
                    _set_run_style(row[4].paragraphs[0].runs[0], size=10, bold=True, color=theme["danger"])
                elif status == "low":
                    _shade_cell(row[4], "3D2F13")
                    _set_run_style(row[4].paragraphs[0].runs[0], size=10, bold=True, color=theme["warning"])
                elif status == "normal":
                    _shade_cell(row[4], "123226")
                    _set_run_style(row[4].paragraphs[0].runs[0], size=10, bold=True, color=theme["success"])
                else:
                    _shade_cell(row[4], theme["panel_soft"])
                    _set_run_style(row[4].paragraphs[0].runs[0], size=10, bold=True, color=theme["muted"])
        else:
            _add_small("No vitals captured.", color="64748B")

        _add_h1("3. Risk Signals")
        risk_scores = analytics.get("risk_scores") or []
        if risk_scores:
            risk_tbl = doc.add_table(rows=1, cols=4)
            risk_tbl.style = "Table Grid"
            for i, txt in enumerate(["Score", "Value", "Band", "Explanation"]):
                c = risk_tbl.rows[0].cells[i]
                c.text = txt
                _shade_cell(c, theme["panel_dark"])
                _set_run_style(c.paragraphs[0].runs[0], size=10, bold=True, color=theme["accent"])
            for r in risk_scores:
                try:
                    pct = int(max(0.0, min(1.0, float(r.get("value", 0.0) or 0.0))) * 100)
                except Exception:
                    pct = 0
                row = risk_tbl.add_row().cells
                row[0].text = str(r.get("name", "risk_score"))
                row[1].text = f"{pct}%"
                row[2].text = str(r.get("band", "low"))
                row[3].text = str(r.get("explanation", ""))
                band = str(r.get("band", "low")).lower()
                if band == "high":
                    _shade_cell(row[2], "3F1C2B")
                    _set_run_style(row[2].paragraphs[0].runs[0], size=10, bold=True, color=theme["danger"])
                elif band == "medium":
                    _shade_cell(row[2], "3D2F13")
                    _set_run_style(row[2].paragraphs[0].runs[0], size=10, bold=True, color=theme["warning"])
                else:
                    _shade_cell(row[2], "123226")
                    _set_run_style(row[2].paragraphs[0].runs[0], size=10, bold=True, color=theme["success"])
        else:
            _add_small("No risk scores captured.", color="64748B")

        _add_h1("4. Diagnosis Support (Possible)")
        diffs = diagnosis.get("differential_diagnoses") or []
        evidence = diagnosis.get("evidence") or []
        missing_questions = diagnosis.get("missing_questions") or []
        doc.add_paragraph("Differential diagnoses (possible):")
        if diffs:
            for d in diffs[:10]:
                doc.add_paragraph(str(d), style="List Bullet")
        else:
            _add_small("No differentials captured.")
        doc.add_paragraph("Evidence mapping:")
        if evidence:
            for e in evidence[:12]:
                doc.add_paragraph(f"{e.get('finding', '')} ({e.get('source', '')})", style="List Bullet")
        else:
            _add_small("No evidence captured.")
        if missing_questions:
            doc.add_paragraph("Missing questions:")
            for q in missing_questions[:12]:
                doc.add_paragraph(str(q), style="List Bullet")

        _add_h1("5. Actionable Insights")
        steps = actionable.get("recommended_next_steps") or {}
        for label in ("Immediate", "Today", "Follow-up"):
            p = doc.add_paragraph()
            r = p.add_run(label)
            _set_run_style(r, size=11, bold=True, color="0F172A")
            for step in steps.get(label) or []:
                doc.add_paragraph(str(step), style="List Bullet")

        gaps = actionable.get("documentation_gaps") or []
        if gaps:
            doc.add_paragraph("Documentation gaps:")
            for g in gaps[:15]:
                doc.add_paragraph(
                    f"[{g.get('severity', 'low')}] {g.get('gap', '')} - {g.get('why_it_matters', '')}",
                    style="List Bullet",
                )

        safety = actionable.get("safety_net") or []
        if safety:
            doc.add_paragraph("Safety net guidance:")
            for s in safety[:15]:
                doc.add_paragraph(str(s), style="List Bullet")

        doc.add_page_break()
        _add_h1("Appendix: Raw Structured Snapshot")
        raw = {
            "chief_complaint": structured_data.get("chief_complaint"),
            "assessment": structured_data.get("assessment"),
            "treatment_plan": structured_data.get("treatment_plan"),
            "red_flags": (analytics.get("red_flag_summary") or {}).get("items") or [],
        }
        raw_p = doc.add_paragraph(json.dumps(raw, indent=2, ensure_ascii=False))
        for run in raw_p.runs:
            _set_run_style(run, size=9, color=theme["muted"])
        doc.save(str(out_path))
        return {
            "type": "docx",
            "name": out_path.name,
            "path": str(out_path),
            "size": int(out_path.stat().st_size),
        }
    except Exception:
        return None


def _generate_with_python_pptx(out_path: Path, structured_data: Dict[str, Any], insights: Dict[str, Any]) -> bool:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return False

    try:
        from pptx.util import Inches, Pt  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
    except Exception:
        return False

    def _title(slide, text: str) -> None:
        t = slide.shapes.title
        if t is not None:
            t.text = text
            if t.text_frame and t.text_frame.paragraphs:
                run = t.text_frame.paragraphs[0].runs[0] if t.text_frame.paragraphs[0].runs else None
                if run:
                    run.font.bold = True
                    run.font.size = Pt(28)
                    run.font.color.rgb = RGBColor(15, 23, 42)

    def _add_box(slide, title: str, body: str, left: float, top: float, width: float, height: float) -> None:
        shp = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = shp.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(14)
        p1.font.color.rgb = RGBColor(15, 23, 42)
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(11)
        p2.font.color.rgb = RGBColor(51, 65, 85)

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        chief = ((structured_data.get("chief_complaint") or {}).get("complaint")) or "not captured"
        risk_tier = ((insights.get("diagnosis_support") or {}).get("risk_tier")) or "low"
        red_count = (((insights.get("analytics") or {}).get("red_flag_summary") or {}).get("count")) or 0
        gap_count = len(((insights.get("actionable_insights") or {}).get("documentation_gaps")) or [])
        conf = structured_data.get("confidence_score", 0.5)
        conf_pct = int(max(0.0, min(1.0, float(conf or 0.0))) * 100)

        s1 = prs.slides.add_slide(prs.slide_layouts[5])
        _title(s1, "Clinical Intelligence Report")
        _add_box(
            s1,
            "Executive Summary",
            f"Chief complaint: {chief}\nRisk tier: {risk_tier}\nRed flags: {red_count}\nDocumentation gaps: {gap_count}\nConfidence: {conf_pct}%",
            0.5,
            1.2,
            12.2,
            4.6,
        )
        _add_box(
            s1,
            "Disclaimer",
            "Decision support only; clinician verification required.\nNot a final diagnosis or prescription.",
            0.5,
            6.0,
            12.2,
            1.0,
        )

        s2 = prs.slides.add_slide(prs.slide_layouts[5])
        _title(s2, "Vitals and Risk")
        vitals = ((insights.get("analytics") or {}).get("vitals")) or []
        risk_scores = ((insights.get("analytics") or {}).get("risk_scores")) or []
        v_lines = ["Vital | Value | Range | Status"]
        for v in vitals[:12]:
            v_lines.append(
                f"{v.get('name', '')} | {v.get('value', 'not captured')} {v.get('unit', '')}".strip()
                + f" | {v.get('min', '-')} - {v.get('max', '-')}"
                + f" | {v.get('status', 'not_captured')}"
            )
        _add_box(s2, "Vitals", "\n".join(v_lines), 0.5, 1.1, 6.2, 5.8)
        r_lines = ["Score | Value | Band"]
        for r in risk_scores[:10]:
            pct = int(max(0.0, min(1.0, float(r.get("value", 0.0) or 0.0))) * 100)
            r_lines.append(f"{r.get('name', '')} | {pct}% | {r.get('band', 'low')}")
        _add_box(s2, "Risk Scores", "\n".join(r_lines), 6.9, 1.1, 5.8, 5.8)

        s3 = prs.slides.add_slide(prs.slide_layouts[5])
        _title(s3, "Diagnosis Support and Actions")
        ds = insights.get("diagnosis_support") or {}
        ai = insights.get("actionable_insights") or {}
        diffs = ds.get("differential_diagnoses") or []
        evidence = ds.get("evidence") or []
        steps = ai.get("recommended_next_steps") or {}
        safety = ai.get("safety_net") or []
        left = "Differentials (possible):\n" + ("\n".join([f"- {d}" for d in diffs[:8]]) or "- none")
        left += "\n\nEvidence:\n" + ("\n".join([f"- {e.get('finding','')} ({e.get('source','')})" for e in evidence[:8]]) or "- none")
        right = "Immediate:\n" + ("\n".join([f"- {x}" for x in (steps.get("Immediate") or [])[:5]]) or "- none")
        right += "\n\nToday:\n" + ("\n".join([f"- {x}" for x in (steps.get("Today") or [])[:5]]) or "- none")
        right += "\n\nFollow-up:\n" + ("\n".join([f"- {x}" for x in (steps.get("Follow-up") or [])[:5]]) or "- none")
        right += "\n\nSafety Net:\n" + ("\n".join([f"- {x}" for x in safety[:6]]) or "- none")
        _add_box(s3, "Clinical Reasoning", left, 0.5, 1.1, 6.2, 5.8)
        _add_box(s3, "Action Plan", right, 6.9, 1.1, 5.8, 5.8)

        prs.save(str(out_path))
        return True
    except Exception:
        return False
